"""
pybaidu.library
百度文库爬虫接口

作者: stripe-python
更新时间: 2022.1.13

实例1(下载文字)：
>>> from pybaidu.library import *
>>> # 这里的URL来源于网络
>>> library = BaiduLibrary(url=(
>>>     'https://wenku.baidu.com/view/'
>>>     'b81a88f992c69ec3d5bbfd0a79563c1ec4dad7f2.html?'
>>>     'fixfr=tgxz9vqleU5q8%252FyFx4xBqQ%253D%253D&fr=income10-wk_app_search_ctrX-search'),
>>>     cookie='你在开发者工具中的cookie')
>>> print(library.get_text())  # 打印文字
>>> library.save_txt('python.txt')  # 保存至本地TXT文件

实例2(下载PPT)：
>>> from pybaidu.library import *
>>> # 这里的URL来源于网络
>>> library = BaiduLibrary(
>>>     url=(
>>>        'https://wenku.baidu.com/view/'
>>>        'c51a1a75b5360b4c2e3f5727a5e9856a5712260b.html?'
>>>        'fixfr=t8pcDkNYCI%252F7qRYCWsbwMQ%253D%253D&fr=income4-wk_app_search_ctrX-search'
>>>     ),
>>>     cookie='你在开发者工具中的cookie'
>>> )
>>> library.save_pptx('python.pptx')  # 保存至本地PPTX文件
"""

import os as _os
import shutil as _shutil

import requests as _requests
import pptx as _pptx  # 下载PPTX的库
from pptx.util import Cm as _cm  # 缩放图片使用cm
from bs4 import BeautifulSoup as _soup

from pybaidu.locals import *


__all__ = ['BaiduLibrary', 'BaiduWenKu']


class BaiduLibrary(object):
    def __init__(self, url: str, cookie: str = '', proxy_ip: str = None, proxy_type: str = HTTPS,
                 encoding: str = 'utf-8', parser=HTML_PARSER):
        """
        百度文库爬虫核心类。
        提供save_txt、save_pptx和get_text函数。
        在未登录或为输入cookie情况下，爬虫信息会不完整。
        
        :param url: 文章URL
        :param cookie: 你在开发者工具中的cookie
        :param proxy_ip: 代理IP
        :param proxy_type: 代理IP方式
        :param encoding: 网页编码，建议UTF-8
        :param parser: BeautifulSoup4解析器，默认html.parser，支持html.parser、lxml、html5lib、xml、lxml-xml等
        :type url: str
        :type cookie: str
        :type proxy_ip: str
        :type proxy_type: str
        :type encoding: str
        :type parser: str
        """
        self.encoding = encoding
        self.url = url
        self.headers = {
            'Cookie': cookie,
            'User-Agent': USER_AGENT
        }
        self.parser = parser
        self._use_proxy_ip = proxy_ip is not None
        self.proxies = {proxy_type: proxy_ip}
        self._request()
        # 定义API
        self.get = self.get_text
        self.save = self.save_txt
    
    def _request(self):
        if self._use_proxy_ip:
            self.response = _requests.get(self.url, headers=self.headers, proxies=self.proxies)
        else:
            self.response = _requests.get(self.url, headers=self.headers)
        self.response.encoding = self.encoding
        self.soup = _soup(self.response.text, self.parser)
    
    def get_text(self):
        """
        获得该网站的文字。
        会将换行或TAB去掉。
        
        :rtype: str
        :return: 文字信息
        """
        texts = self.soup.find_all('p', attrs={'v-pre': '', 'class': 'reader-word-layer'})
        texts = [i.string for i in texts]
        result = ''.join(texts)
        result = result.replace('\u2002', ' ')  # 替换HTML实体半角空格(ensp)
        return result
    
    def save_txt(self, save_path='article.txt', mode='w', encoding='utf-8', **kwargs):
        r"""
        将该网站的文字信息写入本地TXT文件。
        使用了with-open，文件安全性好。
        关键字实参的详细描述参考python的open函数和python文件IO操作文档
        
        :param save_path: 保存的TXT文件路径。
        :param mode: 写入模式，支持w、a、w+，详情参考open函数。
        :param encoding: TXT文件编码，默认UTF-8
        :keyword buffering: 如果被设为 0，则不会有寄存。如果值取 1，访问文件时会寄存行。如果将值设为大于1的整数，则表示寄存区的缓冲大小。如果取负值，寄存区的缓冲大小则为系统默认。
        :keyword errors: 异常处理方式，当取strict时，字符编码出现问题时会报错；当取ignore时，编码出现问题，程序会忽略而过，继续执行下面的程序。
        :keyword newline: 区分换行符，如\n、空字符串、None、\r、\r\n。
        :keyword closefd: 取值为False时，文件必须为文件描述符。
        :keyword opener: 自定义文件打开器。
        :type save_path: str
        :type mode: str
        :type encoding: str
        :rtype: str
        :return: 保存的文件路径
        """
        with open(save_path, mode, encoding=encoding, **kwargs) as fp:
            fp.write(self.get_text())
        return save_path
    
    def save_pptx(self, save_path='article.pptx', progress_bar=False, temp_folder='temp',
                  delete_temp_folder=True, **kwargs):
        """
        下载图片集，并整合为PPTX。
        使用python-pptx模块。
        
        :param save_path: 保存的PPTX文件路径。
        :param progress_bar: 是否显示进度条。
        :param temp_folder: 临时图片集文件夹，建议使用不存在的路径
        :param delete_temp_folder: 是否删除临时图片集文件夹
        :keyword pptx: 临时PPTX文件路径，默认由python-pptx模块指定。
        :rtype: str
        :return: 保存的PPTX文件路径
        """
        if not _os.path.isdir(temp_folder):
            _os.mkdir(temp_folder)
        div_list = self.soup.find_all('div', attrs={'class': 'ppt-image-wrap'})
        img_list = [div.find('img') for div in div_list]  # 列出图片标签
        
        ppt_file = _pptx.Presentation(**kwargs)  # 创建PPTX对象
        
        for n, img in enumerate(img_list):
            if progress_bar:
                print(PROGRESS_BAR, end='', flush=True)
            url = img.get('src', img.get('data-src'))  # 获取图片URL
            if self._use_proxy_ip:
                response = _requests.get(url, headers=self.headers, proxies=self.proxies)
            else:
                response = _requests.get(url, headers=self.headers)
            path = _os.path.join(temp_folder, f'{n}.jpg')  # 下载图片路径
            with open(path, 'wb') as image_file:
                image_file.write(response.content)  # 保存图片
            title_slide_layout = ppt_file.slide_layouts[6]  # 使用第七种版式，即Blank
            slide = ppt_file.slides.add_slide(title_slide_layout)  # 创建新页
            slide.shapes.add_picture(path, 0, 0, _cm(PPT_WIDTH), _cm(PPT_HEIGHT))  # 写入图片并缩放
                
        if delete_temp_folder:
            _shutil.rmtree(temp_folder)  # 删除临时图片集文件夹
    
        ppt_file.save(save_path)
        return save_path
    
    
BaiduWenKu = BaiduLibrary  # 定义API
